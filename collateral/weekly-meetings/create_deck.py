import os
import collections
import collections.abc
try:
    collections.Mapping = collections.abc.Mapping
    collections.Iterable = collections.abc.Iterable
    collections.Sequence = collections.abc.Sequence
except AttributeError: pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    DARK_BG, LIGHT_BG = RGBColor(0x1E, 0x27, 0x61), RGBColor(0xF7, 0xF9, 0xFC)
    ACCENT = RGBColor(0x4F, 0xC3, 0xF7)
    TEXT_WHITE, TEXT_LIGHT_BLUE = RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xCA, 0xDC, 0xFC)
    TEXT_NAVY, TEXT_MUTED = RGBColor(0x1E, 0x27, 0x61), RGBColor(0x5A, 0x6E, 0x8C)
    DANGER, SUCCESS = RGBColor(0xF9, 0x61, 0x67), RGBColor(0x69, 0xF0, 0xAE)
    
    def add_slide(bg=LIGHT_BG):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
        return slide
        
    def txt(slide, text, l, t, w, h, size=14, color=TEXT_NAVY, bold=False, italic=False, align=None, name="Calibri", name_bold="Calibri Bold", name_light="Calibri Light"):
        tb = slide.shapes.add_textbox(l, t, w, h)
        p = tb.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = name_bold if bold else (name_light if size < 18 else name)
        if align: p.alignment = align
        return tb
        
    def rect(slide, l, t, w, h, fill=None, fl=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        if fill: shape.fill.solid(); shape.fill.fore_color.rgb = fill
        else: shape.fill.background()
        if fl: shape.line.color.rgb = fl; shape.line.width = Pt(1)
        else: shape.line.fill.background()
        return shape

    def card(slide, l, t, w, h, acc=ACCENT, bg=TEXT_WHITE):
        rect(slide, l, t, w, h, fill=bg, fl=TEXT_NAVY)
        rect(slide, l, t, Inches(0.1), h, fill=acc, fl=acc)

    def notes(slide, n):
        slide.notes_slide.notes_text_frame.text = n

    def add_divider(num, title, sub, spec_note):
        s = add_slide(DARK_BG)
        # Big watermark
        tb = txt(s, f"{num:02d}", Inches(5), Inches(1), Inches(3), Inches(4), size=120, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        # tb.text_frame.paragraphs[0].font.color.theme_color = None
        # Fake opacity by matching closer to bg
        tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2A, 0x38, 0x82)  
        txt(s, title, Inches(2), Inches(3), Inches(9.33), Inches(1), size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, sub, Inches(2), Inches(4), Inches(9.33), Inches(0.5), size=20, color=TEXT_LIGHT_BLUE, align=PP_ALIGN.CENTER)
        rect(s, Inches(4), Inches(4.8), Inches(5.33), Inches(0.03), fill=ACCENT, fl=ACCENT)
        notes(s, spec_note)

    # SEC 0 - 01
    s1 = add_slide(DARK_BG)
    txt(s1, "XC GRADIENT", Inches(0.5), Inches(0.5), Inches(3), Inches(0.5), size=14, color=ACCENT, bold=True)
    rect(s1, Inches(6.36), Inches(2), Inches(0.6), Inches(0.6), fill=TEXT_WHITE, fl=TEXT_WHITE)
    txt(s1, "Building the Machine", Inches(2), Inches(3), Inches(9.33), Inches(1), size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s1, "XC Gradient Company Operating System — April 7, 2026", Inches(2), Inches(3.8), Inches(9.33), Inches(0.5), size=18, color=TEXT_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    txt(s1, "Oriol · Arnau · Adam", Inches(2), Inches(6.5), Inches(9.33), Inches(0.5), size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    notes(s1, "Good morning. Today's meeting has one job: by the time we leave this room, we have a system. A real operating system for this company. Not slides that sit in a Drive folder — an actual running infrastructure that we use starting today. Give me 30 minutes.")

    # 02
    s2 = add_slide(DARK_BG)
    txt(s2, "01 — Notion OS\n02 — OKRs & Phase Model\n03 — Daily Logging System\n04 — Discord\n05 — Google Drive\n06 — Meeting Transcription", Inches(1.5), Inches(2), Inches(6), Inches(4), size=18, color=TEXT_WHITE, name="Calibri Bold")
    for i in range(6): txt(s2, f"0{i+1}", Inches(1.5), Inches(2 + (i*0.5)), Inches(1), Inches(0.5), size=28, color=ACCENT, bold=True)
    for i in range(6): rect(s2, Inches(8), Inches(2.2 + (i*0.5)), Inches(0.2), Inches(0.2), fill=ACCENT, fl=ACCENT)
    rect(s2, Inches(8.09), Inches(2.3), Inches(0.02), Inches(2.5), fill=TEXT_MUTED, fl=TEXT_MUTED)
    txt(s2, "Duration: ~30 min", Inches(1.5), Inches(6.5), Inches(3), Inches(0.5), size=12, color=TEXT_MUTED)
    notes(s2, "Six topics. Some will take 5 minutes, some 10. I'll tell you upfront: by the end of this, I need each of you to fill in your Q2 personal KRs live in Notion. That's the deliverable of this meeting.")

    # SEC 1 - 03
    add_divider(1, "Our Company Brain", "Notion as the single source of truth", "Notion is not a note-taking app. It's the company. If it's not in Notion, it didn't happen.")

    # 04
    s4 = add_slide()
    card(s4, Inches(0.5), Inches(1), Inches(6), Inches(1.5))
    txt(s4, "Document decisions, not activity", Inches(0.8), Inches(1.2), Inches(5), Inches(0.5), size=16, bold=True)
    txt(s4, "If it takes more than 2 minutes to log, it won't get logged. We only write what a future cofounder or investor would need to know.", Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.8), size=14, color=TEXT_MUTED)
    
    card(s4, Inches(0.5), Inches(3), Inches(6), Inches(1.5))
    txt(s4, "Passive awareness", Inches(0.8), Inches(3.2), Inches(5), Inches(0.5), size=16, bold=True)
    txt(s4, "Notion is open all day. In 30 seconds you know where each person is, what's blocked, and what's at risk this week. Without sending a message.", Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.8), size=14, color=TEXT_MUTED)
    
    card(s4, Inches(0.5), Inches(5), Inches(6), Inches(1.5))
    txt(s4, "Oriol owns the structure", Inches(0.8), Inches(5.2), Inches(5), Inches(0.5), size=16, bold=True)
    txt(s4, "Arnau and Adam never touch Notion architecture. Oriol extracts, writes, and distributes. Your cognitive resources go entirely to building and selling.", Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.8), size=14, color=TEXT_MUTED)

    rect(s4, Inches(8.5), Inches(2), Inches(3), Inches(1), fill=RGBColor(0x3,0x9,0x3B), fl=TEXT_WHITE)
    txt(s4, "Drive", Inches(8.5), Inches(2.2), Inches(3), Inches(1), size=18, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    rect(s4, Inches(8), Inches(3.5), Inches(4), Inches(1), fill=DARK_BG, fl=TEXT_WHITE)
    txt(s4, "Discord", Inches(8), Inches(3.7), Inches(4), Inches(1), size=18, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    rect(s4, Inches(7.5), Inches(5), Inches(5), Inches(1), fill=RGBColor(0x2A,0x38,0x82), fl=TEXT_WHITE)
    txt(s4, "Notion", Inches(7.5), Inches(5.2), Inches(5), Inches(1), size=18, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    notes(s4, "The rule is simple: Discord is what we're thinking about right now. Notion is what we decided. Drive is for files too big for Notion. That's the whole system.")

    # 05
    s5 = add_slide()
    txt(s5, "The six sections", Inches(0.5), Inches(0.5), Inches(5), Inches(0.5), size=22, bold=True)
    for i, (ic, t, d) in enumerate([("🏠", "Home", "Your daily landing page"), ("🎯", "Goals & OKRs", "Full goal cascade"), ("📓", "Daily Logs", "Oriol / Arnau / Adam"), ("📅", "Weekly Syncs", "Meeting outputs + AI summaries"), ("🧠", "Company Brain", "Strategy, decisions, pipeline"), ("📁", "Archive", "Anything older than 2 quarters")]):
        r, c = i // 3, i % 3
        card(s5, Inches(0.5 + c*4.2), Inches(2 + r*2), Inches(3.8), Inches(1.5))
        txt(s5, f"{ic}  {t}", Inches(0.8 + c*4.2), Inches(2.2 + r*2), Inches(3.3), Inches(0.5), size=18, bold=True)
        txt(s5, d, Inches(0.8 + c*4.2), Inches(2.7 + r*2), Inches(3.3), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s5, "One rule: if it's not in one of these six, it doesn't exist.", Inches(0.5), Inches(6.5), Inches(12), Inches(0.5), size=13, color=TEXT_NAVY, italic=True)
    rect(s5, Inches(0.5), Inches(6.8), Inches(12), Inches(0.02), fill=ACCENT, fl=ACCENT)
    notes(s5, "This is the entire company knowledge base. Six sections. I'll show you each one. The most important one is Company Brain — specifically the Decision Log. That's the one we'll be most disciplined about.")

    # 06
    s6 = add_slide()
    txt(s6, "1  Phase Banner", Inches(0.5), Inches(1), Inches(4), Inches(0.5), size=16, bold=True)
    txt(s6, "Current phase + north star metric. Updated manually at phase gates.", Inches(0.5), Inches(1.4), Inches(4), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s6, "2  Big 3 This Week", Inches(0.5), Inches(2.2), Inches(4), Inches(0.5), size=16, bold=True)
    txt(s6, "Three numbers that tell us if the company is alive and growing.", Inches(0.5), Inches(2.6), Inches(4), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s6, "3  Streak Tracker", Inches(0.5), Inches(3.4), Inches(4), Inches(0.5), size=16, bold=True)
    txt(s6, "Daily log streak per person. Resets to zero if you miss a day.", Inches(0.5), Inches(3.8), Inches(4), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s6, "4  Company Calendar", Inches(0.5), Inches(4.6), Inches(4), Inches(0.5), size=16, bold=True)
    txt(s6, "2-week rolling view. Only company-wide events.", Inches(0.5), Inches(5.0), Inches(4), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s6, "5  Quick Links", Inches(0.5), Inches(5.8), Inches(4), Inches(0.5), size=16, bold=True)
    txt(s6, "Sprint card, Decision Log, PoC Pipeline, latest sync.", Inches(0.5), Inches(6.2), Inches(4), Inches(0.5), size=14, color=TEXT_MUTED)
    
    # Mockup
    rect(s6, Inches(5.5), Inches(0.5), Inches(7.33), Inches(6.5), fill=RGBColor(0xEE, 0xEE, 0xEE), fl=TEXT_MUTED)
    rect(s6, Inches(5.8), Inches(1), Inches(6.73), Inches(0.8), fill=DARK_BG)
    txt(s6, "Phase 1: PoC Execution | Metric: # PoCs", Inches(6), Inches(1.2), Inches(6), Inches(0.5), size=16, color=TEXT_WHITE, bold=True)
    for i in range(3): rect(s6, Inches(5.8 + i*2.3), Inches(2), Inches(2.1), Inches(1), fill=TEXT_WHITE, fl=TEXT_MUTED)
    rect(s6, Inches(5.8), Inches(3.2), Inches(3.3), Inches(1.5), fill=TEXT_WHITE, fl=TEXT_MUTED) # streak
    txt(s6, "Streaks: Oriol (1) Arnau (0) Adam (0)", Inches(6), Inches(3.4), Inches(3), Inches(1), size=12)
    rect(s6, Inches(9.3), Inches(3.2), Inches(3.2), Inches(1.5), fill=TEXT_WHITE, fl=TEXT_MUTED) # cal
    rect(s6, Inches(5.8), Inches(4.9), Inches(6.73), Inches(1.5), fill=TEXT_WHITE, fl=TEXT_MUTED) # links
    notes(s6, "This is what you see when you open your laptop. Before you open anything else — terminal, IDE, whatever — Notion is open. The Home page tells you in 30 seconds if anything needs your attention.")

    # 07
    s7 = add_slide()
    txt(s7, "Four tools. One job each.", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    def tcard(l, t, name, role, rule, acc):
        card(s7, l, t, Inches(5), Inches(2), acc=acc)
        txt(s7, name, l + Inches(0.3), t + Inches(0.2), Inches(4.5), Inches(0.5), size=20, bold=True)
        txt(s7, role, l + Inches(0.3), t + Inches(0.7), Inches(4.5), Inches(0.5), size=14, color=TEXT_MUTED)
        txt(s7, rule, l + Inches(0.3), t + Inches(1.2), Inches(4.5), Inches(0.5), size=14, bold=True)

    tcard(Inches(1.5), Inches(1.5), "Notion", "Company OS — all structured knowledge", "If it's a decision, goal, log, or summary → Notion", ACCENT)
    tcard(Inches(7), Inches(1.5), "Discord", "Real-time async comms + voice", "Working memory. Ephemeral. Decisions migrate to Notion.", RGBColor(0x58,0x65,0xF2))
    tcard(Inches(1.5), Inches(4), "Drive", "Binary blobs only", "PPTs, contracts, invoices, CSV exports, large files", RGBColor(0x34,0xA8,0x53))
    tcard(Inches(7), Inches(4), "WhatsApp", "Personal / emergency only", "Urgent pings that cannot wait for Discord", RGBColor(0x25,0xD3,0x66))
    
    card(s7, Inches(1.5), Inches(6.5), Inches(10.5), Inches(0.6), acc=DANGER, bg=RGBColor(0xFD, 0xE2, 0xE3))
    txt(s7, "❌ Slack: not used. Redundant with Discord at 3 people. Revisit at first hire.", Inches(1.8), Inches(6.6), Inches(10), Inches(0.5), size=14, color=TEXT_NAVY, bold=True)
    notes(s7, "Slack is a deliberate no. We don't need three messaging systems. When we hire the first employee, we can revisit — they may expect Slack for HR tools. That's a 2027 problem.")

    # SEC 2 - 08
    add_divider(2, "Where We're Going", "Phase model, growth math, and our 2026 OKR cascade", "This is the strategic layer. I want you to understand not just what the goals are, but why they're structured this way.")

    # 09
    s9 = add_slide()
    txt(s9, "Two different questions", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    # KPI Card
    card(s9, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), acc=DARK_BG)
    rect(s9, Inches(0.8), Inches(1.8), Inches(0.8), Inches(0.3), fill=DARK_BG)
    txt(s9, "KPI", Inches(0.8), Inches(1.8), Inches(0.8), Inches(0.3), size=14, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s9, "Is the business healthy right now?", Inches(0.8), Inches(2.3), Inches(5), Inches(0.5), size=18, bold=True)
    txt(s9, "Permanent dials on your dashboard. Never disappear. Measured every week forever.", Inches(0.8), Inches(2.9), Inches(5), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s9, "• ARR\n• Pipeline count\n• Burn rate\n• A score\n• Runway", Inches(1), Inches(3.8), Inches(4), Inches(1.5), size=14, color=TEXT_NAVY)
    txt(s9, "READ WEEKLY", Inches(0.8), Inches(6.5), Inches(3), Inches(0.3), size=11, color=TEXT_MUTED, bold=True)
    
    # OKR Card
    card(s9, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5), acc=ACCENT)
    rect(s9, Inches(7.3), Inches(1.8), Inches(0.8), Inches(0.3), fill=ACCENT)
    txt(s9, "OKR", Inches(7.3), Inches(1.8), Inches(0.8), Inches(0.3), size=14, color=TEXT_NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(s9, "Where are we trying to go this quarter?", Inches(7.3), Inches(2.3), Inches(5), Inches(0.5), size=18, bold=True)
    txt(s9, "Temporary targets with an expiry date. Die at end of quarter. Hit or missed, then replaced.", Inches(7.3), Inches(2.9), Inches(5), Inches(0.5), size=14, color=TEXT_MUTED)
    txt(s9, "Objective: Prove product-market fit in industrial SMEs\n   KR1: 2 paying PoCs by Jun 30\n   KR2: A score > 0.80 on Decfa corpus\n   KR3: 5 qualified pipeline conversations", Inches(7.5), Inches(3.8), Inches(5), Inches(1.5), size=14, color=TEXT_NAVY)
    txt(s9, "REVIEWED QUARTERLY", Inches(7.3), Inches(6.5), Inches(3), Inches(0.3), size=11, color=TEXT_MUTED, bold=True)
    
    txt(s9, "The relationship: OKR defines the destination. KPI tells you if you're moving toward it.", Inches(0.5), Inches(7.2), Inches(12), Inches(0.3), size=14, color=TEXT_NAVY, italic=True, align=PP_ALIGN.CENTER)
    notes(s9, "You can have great KPIs and be building the wrong thing. OKRs answer 'what matters most right now?' — and that answer changes every 13 weeks.")

    # 10
    s10 = add_slide(DARK_BG)
    txt(s10, "The ARR equation", Inches(0.5), Inches(0.5), Inches(12), Inches(0.8), size=36, color=TEXT_WHITE, bold=True)
    
    txt(s10, "Annual target", Inches(0.5), Inches(2), Inches(5), Inches(0.3), size=13, color=TEXT_MUTED)
    txt(s10, "ARR × 3", Inches(0.5), Inches(2.3), Inches(5), Inches(0.5), size=22, color=TEXT_WHITE, bold=True)
    
    txt(s10, "Quarterly target", Inches(0.5), Inches(3), Inches(5), Inches(0.3), size=13, color=TEXT_MUTED)
    txt(s10, "ARR × 1.75 per quarter", Inches(0.5), Inches(3.3), Inches(5), Inches(0.5), size=22, color=TEXT_WHITE, bold=True)
    
    txt(s10, "Weekly target", Inches(0.5), Inches(4), Inches(5), Inches(0.3), size=13, color=TEXT_MUTED)
    txt(s10, "ARR × 1.06 per week", Inches(0.5), Inches(4.3), Inches(5), Inches(0.5), size=22, color=TEXT_WHITE, bold=True)
    
    txt(s10, "= +6% week over week", Inches(0.5), Inches(5.2), Inches(6), Inches(0.8), size=32, color=ACCENT, bold=True)
    
    rect(s10, Inches(7), Inches(2), Inches(4), Inches(2), fill=DARK_BG, fl=TEXT_WHITE)
    txt(s10, "5–7%", Inches(7), Inches(2.2), Inches(4), Inches(1), size=64, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s10, "WoW growth = good", Inches(7), Inches(3.2), Inches(4), Inches(0.5), size=14, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    rect(s10, Inches(7), Inches(4.2), Inches(4), Inches(2), fill=ACCENT, fl=ACCENT)
    txt(s10, "10%", Inches(7), Inches(4.4), Inches(4), Inches(1), size=64, color=TEXT_NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(s10, "WoW growth = magic", Inches(7), Inches(5.4), Inches(4), Inches(0.5), size=14, color=TEXT_NAVY, bold=True, align=PP_ALIGN.CENTER)
    
    txt(s10, "Two levers only: N_clients (new logos) and ACV (upsell existing). Every week the question is one of these two — never both.", Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), size=13, color=TEXT_MUTED)
    notes(s10, "At 10% week-over-week, you 5× in a year. At 5%, you 12× in a year. The math is violent once you have any base at all. Right now we have no base. That changes after today.")

    # 11
    s11 = add_slide()
    # Timeline
    rect(s11, Inches(0.5), Inches(0.5), Inches(2.8), Inches(0.5), fill=ACCENT, fl=ACCENT)
    txt(s11, "Phase 0 · Ignition", Inches(0.5), Inches(0.6), Inches(2.8), Inches(0.5), size=12, color=TEXT_NAVY, bold=True, align=PP_ALIGN.CENTER)
    for i, t in enumerate(["Phase 1 · PoC Execution", "Phase 2 · Pipeline", "Phase 3 · Launch"]):
        rect(s11, Inches(3.6 + i*3.1), Inches(0.5), Inches(2.8), Inches(0.5), fill=TEXT_WHITE, fl=TEXT_NAVY)
        txt(s11, t, Inches(3.6 + i*3.1), Inches(0.6), Inches(2.8), Inches(0.5), size=12, color=TEXT_NAVY, align=PP_ALIGN.CENTER)
        txt(s11, "→", Inches(3.3 + i*3.1), Inches(0.6), Inches(0.3), Inches(0.5), size=16, color=TEXT_NAVY, align=PP_ALIGN.CENTER)
        
    card(s11, Inches(2), Inches(1.5), Inches(9.33), Inches(4))
    txt(s11, "📍 Phase 0 — Ignition", Inches(2.3), Inches(1.8), Inches(8), Inches(0.5), size=24, bold=True)
    txt(s11, "Apr 3 – Apr 7, 2026", Inches(2.3), Inches(2.3), Inches(8), Inches(0.3), size=14, color=TEXT_MUTED)
    txt(s11, "Is the team operationally ready to execute?", Inches(2.3), Inches(2.8), Inches(8), Inches(0.5), size=18, italic=True)
    
    txt(s11, "North Star: Demo readiness + OS installed — binary by April 7\n\nPhase gate: This meeting. Exit criteria: all 3 cofounders in Notion with Q2 KRs authored.\n\nYour USP: Fully on-prem RAG — zero cloud dependency, NIS2-native", Inches(2.5), Inches(3.5), Inches(8.5), Inches(1.5), size=14)
    
    card(s11, Inches(2), Inches(6), Inches(9.33), Inches(0.8), acc=DANGER, bg=RGBColor(0xFD, 0xE2, 0xE3))
    txt(s11, "⚡ This afternoon: Paver visit. Three possible outcomes → all are fine. Any result gets logged.", Inches(2.3), Inches(6.2), Inches(8), Inches(0.5), size=13, color=TEXT_NAVY)
    notes(s11, "Phase 0 is the only phase that's about us, not the market. The question is: are we operationally ready? By end of this meeting, yes.")

    # 12
    s12 = add_slide()
    rect(s12, Inches(0.5), Inches(0.5), Inches(2.8), Inches(0.5), fill=TEXT_WHITE, fl=TEXT_NAVY)
    txt(s12, "Phase 0 · Ignition", Inches(0.5), Inches(0.6), Inches(2.8), Inches(0.5), size=12, color=TEXT_NAVY, align=PP_ALIGN.CENTER)
    txt(s12, "→", Inches(3.3), Inches(0.6), Inches(0.3), Inches(0.5), size=16, color=TEXT_NAVY, align=PP_ALIGN.CENTER)
    rect(s12, Inches(3.6), Inches(0.5), Inches(2.8), Inches(0.5), fill=ACCENT, fl=ACCENT)
    txt(s12, "Phase 1 · PoC Execution", Inches(3.6), Inches(0.6), Inches(2.8), Inches(0.5), size=12, color=TEXT_NAVY, bold=True, align=PP_ALIGN.CENTER)
    for i, t in enumerate(["Phase 2 · Pipeline", "Phase 3 · Launch"]):
        rect(s12, Inches(6.7 + i*3.1), Inches(0.5), Inches(2.8), Inches(0.5), fill=TEXT_WHITE, fl=TEXT_NAVY)
        txt(s12, t, Inches(6.7 + i*3.1), Inches(0.6), Inches(2.8), Inches(0.5), size=12, color=TEXT_NAVY, align=PP_ALIGN.CENTER)
        txt(s12, "→", Inches(6.4 + i*3.1), Inches(0.6), Inches(0.3), Inches(0.5), size=16, color=TEXT_NAVY, align=PP_ALIGN.CENTER)

    card(s12, Inches(1), Inches(1.5), Inches(11.3), Inches(4.2))
    txt(s12, "Phase 1 — PoC Execution", Inches(1.3), Inches(1.8), Inches(8), Inches(0.5), size=24, bold=True)
    txt(s12, "Apr 7 → Jul 1, 2026", Inches(1.3), Inches(2.3), Inches(8), Inches(0.3), size=14, color=TEXT_MUTED)
    txt(s12, "Can we produce a result a client would pay for?", Inches(1.3), Inches(2.7), Inches(8), Inches(0.5), size=18, italic=True)
    txt(s12, "North Star: # of PoCs with documented, measurable outcome", Inches(1.3), Inches(3.3), Inches(8), Inches(0.5), size=14, bold=True)
    
    txt(s12, "#1  PoC count with A score logged             0 / 2", Inches(1.5), Inches(3.8), Inches(6), Inches(0.3), size=14)
    txt(s12, "#2  Pipeline conversations / week             0", Inches(1.5), Inches(4.2), Inches(6), Inches(0.3), size=14)
    txt(s12, "#3  A score on Decfa corpus                   0.00", Inches(1.5), Inches(4.6), Inches(6), Inches(0.3), size=14)
    
    txt(s12, "Target: 5–7% WoW growth on pipeline conversations. A > 0.80 on Decfa.", Inches(1.3), Inches(5.1), Inches(10), Inches(0.3), size=14)
    
    rect(s12, Inches(1), Inches(6), Inches(11.3), Inches(0.8), fill=ACCENT, fl=ACCENT)
    txt(s12, "Gate: July 1 — Thesis defense. 2 PoCs with hard numbers → pivot from builder to seller.", Inches(1.3), Inches(6.2), Inches(10), Inches(0.5), size=14, color=TEXT_NAVY, bold=True)
    notes(s12, "The thesis defense date is not just academic. July 1 is the day we stop being primarily builders and become primarily sellers. The thesis is our credibility asset — a defended academic RAG paper plus two PoCs with numbers is our pitch deck.")

    # 13
    s13 = add_slide()
    txt(s13, "The back half of 2026", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    rect(s13, Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.5), fill=TEXT_WHITE, fl=TEXT_NAVY)
    txt(s13, "Phase 2 — Pipeline Building", Inches(0.8), Inches(1.8), Inches(5), Inches(0.5), size=18, bold=True)
    txt(s13, "Jul 1 → Sep/Oct", Inches(0.8), Inches(2.2), Inches(5), Inches(0.3), size=13, color=TEXT_MUTED)
    txt(s13, "North Star: # of LOIs / verbal Pioneer commitments (€199/mo)", Inches(0.8), Inches(2.7), Inches(5.2), Inches(0.5), size=14)
    txt(s13, "• #1 LOI count (target: 10)\n• #2 Demos given per week (5–7% WoW)\n• #3 Demo → LOI conversion", Inches(1), Inches(3.3), Inches(5), Inches(1), size=14)
    txt(s13, "Gate: 10 verbal Pioneers before constitution", Inches(0.8), Inches(4.5), Inches(5), Inches(0.3), size=14, bold=True)
    txt(s13, "YC insight: Charging before you can charge. A founder who said yes twice is far more likely to convert than a cold lead.", Inches(0.8), Inches(5.0), Inches(5.2), Inches(0.5), size=12, color=TEXT_MUTED, italic=True)

    rect(s13, Inches(7), Inches(1.5), Inches(5.8), Inches(4.5), fill=TEXT_WHITE, fl=TEXT_NAVY)
    txt(s13, "Phase 3 — Constitution + Launch", Inches(7.3), Inches(1.8), Inches(5), Inches(0.5), size=18, bold=True)
    txt(s13, "Oct → Dec (Q4)", Inches(7.3), Inches(2.2), Inches(5), Inches(0.3), size=13, color=TEXT_MUTED)
    txt(s13, "North Star: ARR (€/month, WoW %)", Inches(7.3), Inches(2.7), Inches(5.2), Inches(0.5), size=14)
    txt(s13, "Launch state: 10 × €199 = €1,990 MRR on day 1 of legal existence", Inches(7.3), Inches(3.2), Inches(5.2), Inches(0.5), size=14)
    txt(s13, "• #1 ARR (5–7% WoW target)\n• #2 Active users per PoC\n• #3 Expansion pipeline", Inches(7.5), Inches(3.7), Inches(5), Inches(0.8), size=14)
    txt(s13, "Gate: End Q4 with first upsell conversation opened", Inches(7.3), Inches(4.7), Inches(5), Inches(0.5), size=14, bold=True)

    rect(s13, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.6), fill=DARK_BG, fl=DARK_BG)
    txt(s13, "Yearly goal: Constitute XC Gradient with 10 paying Pioneer clients by Q4 2026", Inches(1), Inches(6.6), Inches(11.3), Inches(0.5), size=14, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    notes(s13, "Phase 2 is the most underestimated phase. You'll have a working product, a thesis, and two PoCs. You'll feel like the hard part is done. It is not. The hard part is the next 18 months of selling.")

    # 14
    s14 = add_slide()
    txt(s14, "Every task traces up to this", Inches(0.5), Inches(0.5), Inches(8), Inches(0.5), size=22, bold=True)
    
    t = """🏆  YEARLY GOAL
    Constitute XC Gradient with 10 paying Pioneers by Q4 2026
    │
    ├── Q2 2026  (CURRENT — Apr–Jun)
    │   Company: 2 PoCs A>0.80 · 5 qualified leads · Thesis Jul 1
    │   │
    │   ├── ORIOL    KR1: Notion OS live Apr 7
    │   │            KR2: 5 pipeline conversations
    │   │            KR3: Thesis defended Jul 1
    │   │
    │   ├── ARNAU    KR1: 
    │   │            KR2: 
    │   │            KR3: 
    │   │
    │   └── ADAM     KR1: 
    │                KR2: 
    │                KR3: 
    │
    ├── Q3 2026  (stub)   "Build 10 Pioneer LOIs"
    └── Q4 2026  (stub)   "Launch constituted company"
"""
    tb = txt(s14, t, Inches(0.5), Inches(1.2), Inches(8), Inches(6), size=14, name="Courier New")
    
    card(s14, Inches(9), Inches(1.2), Inches(3.8), Inches(1), acc=ACCENT)
    txt(s14, "Arnau and Adam's KRs are blank. We fill them together in the next 15 minutes.", Inches(9.3), Inches(1.4), Inches(3.4), Inches(0.8), size=13, italic=True)
    notes(s14, "The rule on KRs: you define your own. I will never assign KRs to you. Owned KRs get hit. Assigned KRs get ignored. So let's open Notion right now.\n[PAUSE — open Notion live, navigate to Goals & OKRs, build Arnau and Adam's KRs together]")

    # SEC 3 - 15
    add_divider(3, "How We Work Every Day", "The daily log, sprint cards, and streak system", "This is the section most people skip and then regret. The logging system is not bureaucracy. It's the mechanism by which I never have to ask you what you're working on.")

    # 16
    s16 = add_slide()
    txt(s16, "5 minutes. Every day. No exceptions.", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    # Left column rules
    def rule(l, t, icon_color, title, desc):
        rect(s16, l, t, Inches(0.4), Inches(0.4), fill=icon_color, fl=icon_color)
        txt(s16, title, l + Inches(0.6), t - Inches(0.05), Inches(4.5), Inches(0.3), size=16, bold=True)
        txt(s16, desc, l + Inches(0.6), t + Inches(0.3), Inches(4.5), Inches(0.8), size=14, color=TEXT_MUTED)

    rule(Inches(0.5), Inches(1.5), SUCCESS, "Fill before closing Notion", "If you worked today, you log today. Missing a day resets your streak to zero.")
    rule(Inches(0.5), Inches(3.0), ACCENT, "TOMORROW max 3 items", "Forces prioritization. If you have 10 things tomorrow, you don't have a plan — you have a list.")
    rule(Inches(0.5), Inches(4.5), DANGER, "BLOCKERS must be visible", "If something is blocking you, the other two need to see it. That's the whole point of the system.")

    # Right column Mockup
    card(s16, Inches(6.5), Inches(1.5), Inches(6), Inches(5), acc=DARK_BG, bg=RGBColor(0xF0, 0xF2, 0xF5))
    txt(s16, "📅  2026-04-07  ·  Oriol\n🔗  Sprint: W15", Inches(6.8), Inches(1.8), Inches(5), Inches(0.5), size=14, bold=True)
    txt(s16, "✅  DONE TODAY\n    · Pain point map finalized\n    · Paver visit — outcome: qualified lead\n    · Logged 2 decisions", Inches(6.8), Inches(2.5), Inches(5), Inches(0.8), size=14)
    txt(s16, "⏭️  TOMORROW\n    · W15 sprint card with Arnau + Adam\n    · Decfa demo prep check", Inches(6.8), Inches(3.5), Inches(5), Inches(0.6), size=14)
    txt(s16, "⚡  BLOCKERS\n    · Need Arnau: Decfa corpus loaded by Thu?", Inches(6.8), Inches(4.3), Inches(5), Inches(0.4), size=14)
    txt(s16, "📊  GOAL PROGRESS\n    · Notion OS  80% → 100% ✅", Inches(6.8), Inches(5.0), Inches(5), Inches(0.4), size=14)
    notes(s16, "This is what I wrote last night. It took 4 minutes. Arnau, you can see without messaging me that the Paver visit happened, what came out of it, and that I need an answer from you by Thursday. That's the whole value of the system.")

    # 17
    s17 = add_slide()
    txt(s17, "The accountability layer", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    # Left Mockup
    card(s17, Inches(0.5), Inches(1.5), Inches(6.5), Inches(4.5))
    txt(s17, "W15 — Apr 7–13 · Arnau", Inches(0.8), Inches(1.8), Inches(5), Inches(0.3), size=14, bold=True)
    
    txt(s17, "RAG running on Decfa          0%    🔴 Not started\nNetwork topology            100%    ✅ Done\nLatency benchmark < 3s        0%    🔴 Not started", Inches(0.8), Inches(2.3), Inches(6), Inches(1), size=14)
    
    # Progress bars
    rect(s17, Inches(0.8), Inches(2.6), Inches(5), Inches(0.05), fill=DANGER, fl=DANGER)
    rect(s17, Inches(0.8), Inches(3.0), Inches(5), Inches(0.05), fill=SUCCESS, fl=SUCCESS)
    rect(s17, Inches(0.8), Inches(3.4), Inches(5), Inches(0.05), fill=DANGER, fl=DANGER)
    
    txt(s17, "Week outcome (filled Friday EOD): ___", Inches(0.8), Inches(4.5), Inches(5), Inches(0.5), size=14)

    # Right Streak Tracker
    card(s17, Inches(7.5), Inches(1.5), Inches(5.3), Inches(4.5))
    txt(s17, "Daily log streaks", Inches(7.8), Inches(1.8), Inches(4), Inches(0.3), size=14, bold=True)
    txt(s17, "Oriol   🔥  0 days   Best: 0\nArnau   🔥  0 days   Best: 0\nAdam    🔥  0 days   Best: 0", Inches(7.8), Inches(2.3), Inches(4), Inches(1), size=14)
    txt(s17, "⚠️ Warning appears at 7pm if no log submitted today. Missing a day → streak resets to zero. Weekends count if submitted.", Inches(7.8), Inches(4.5), Inches(4.5), Inches(0.8), size=12, color=TEXT_MUTED)

    rect(s17, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), fill=TEXT_WHITE, fl=ACCENT)
    txt(s17, "Gamification rule: we compete on discipline, not output quality. No ranking by performance. Perverse incentives by design avoided.", Inches(0.8), Inches(6.6), Inches(11), Inches(0.5), size=13, italic=True)
    notes(s17, "The streak is the only competitive element. We don't rank each other by how much we shipped. That would create terrible incentives. We rank by whether we showed up and logged. That's all.")

    # 18
    s18 = add_slide()
    txt(s18, "Define nearer things more completely", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    txt(s18, "Top-down cascading — you can only define a smaller goal correctly once you know the bigger goal it serves", Inches(0.5), Inches(1.0), Inches(12), Inches(0.5), size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Table
    table_l, table_t = Inches(1), Inches(1.8)
    rect(s18, table_l, table_t, Inches(11.3), Inches(0.6), fill=DARK_BG, fl=DARK_BG)
    txt(s18, "Cadence                 Duration    Who       Purpose", table_l + Inches(0.5), table_t + Inches(0.1), Inches(10), Inches(0.5), size=14, color=TEXT_WHITE, bold=True)
    rows = [
        ("Daily EOD log", "5 min", "Solo", "Log done / tomorrow / blockers"),
        ("Weekly sync", "60 min", "All 3", "Present week, align sprint cards"),
        ("Monthly close", "90 min", "All 3", "Review month, define next month"),
        ("Quarterly OKR", "Half day", "All 3", "Review Q, define next Q + 3 months"),
        ("Annual planning", "Full day", "All 3", "Yearly goal, stub all 4 quarters")
    ]
    for i, (c1, c2, c3, c4) in enumerate(rows):
        y = table_t + Inches(0.6) + Inches(i * 0.7)
        if i % 2 == 1:
            rect(s18, table_l, y, Inches(11.3), Inches(0.7), fill=RGBColor(0xEE, 0xF1, 0xF6), fl=RGBColor(0xEE, 0xF1, 0xF6))
        txt(s18, f"{c1}", table_l + Inches(0.5), y + Inches(0.2), Inches(2), Inches(0.5), size=14)
        txt(s18, f"{c2}", table_l + Inches(3), y + Inches(0.2), Inches(1.5), Inches(0.5), size=14)
        txt(s18, f"{c3}", table_l + Inches(4.5), y + Inches(0.2), Inches(1.5), Inches(0.5), size=14)
        txt(s18, f"{c4}", table_l + Inches(6), y + Inches(0.2), Inches(4.5), Inches(0.5), size=14)

    card(s18, Inches(1), Inches(6.2), Inches(11.3), Inches(0.8), acc=ACCENT)
    txt(s18, "Q3 and Q4 are always stubs. Detailing Q4 in April is wasted planning — the world will have changed by then.", Inches(1.3), Inches(6.5), Inches(10), Inches(0.5), size=14, color=TEXT_NAVY, italic=True)
    notes(s18, "The most common planning mistake is trying to detail everything at once. We detail only what's next. Everything else is a direction stub that we fill when we get there.")

    # SEC 4 - 19
    s19 = add_slide(DARK_BG)
    tb = txt(s19, "04", Inches(5), Inches(1), Inches(3), Inches(4), size=120, color=RGBColor(0x58,0x65,0xF2), bold=True, align=PP_ALIGN.CENTER)
    # tb.text_frame.paragraphs[0].font.color.theme_color = None
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x38, 0x42, 0x98) # fake opacity
    txt(s19, "Our Communication Layer", Inches(2), Inches(3), Inches(9.33), Inches(1), size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s19, "Discord is working memory. Notion is long-term memory.", Inches(2), Inches(4), Inches(9.33), Inches(0.5), size=20, color=TEXT_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    rect(s19, Inches(4), Inches(4.8), Inches(5.33), Inches(0.03), fill=RGBColor(0x58,0x65,0xF2), fl=RGBColor(0x58,0x65,0xF2))
    notes(s19, "Discord is not for decisions. It's for conversations that lead to decisions. Once something resolves, it goes to Notion. That's the contract.")

    # 20
    s20 = add_slide()
    txt(s20, "One channel, one job", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    # Left Discord channel map
    rect(s20, Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5), fill=RGBColor(0x2B, 0x2D, 0x31), fl=RGBColor(0x2B, 0x2D, 0x31))
    discord_t = "XC GRADIENT\n\n# 📢 announcements    ← decisions\n# 🚨 blockers         ← 24h SLA\n# 🛠️ product-cto      ← Arnau\n# 📊 gtm-coo          ← Adam\n# 💰 finance-cfo      ← Oriol\n# 🔗 resources        ← links\n# 💬 general          ← misc\n\nVOICE\n🔊 The War Room       ← syncs\n🎧 Deep Work          ← silent"
    txt(s20, discord_t, Inches(0.8), Inches(1.8), Inches(4), Inches(4), size=14, color=TEXT_WHITE, name="Courier New")
    
    # Right image with annotations
    image_path = "/mnt/user-data/uploads/1775550409325_image.png"
    if os.path.exists(image_path):
        s20.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(5.5), height=Inches(3.5))

        def callout(text, x, y, w, h, to_x, to_y):
            rect(s20, x, y, w, h, fill=TEXT_WHITE, fl=ACCENT)
            txt(s20, text, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), h - Inches(0.1), size=11, color=TEXT_NAVY, bold=True)
            line = s20.shapes.add_shape(MSO_SHAPE.LINE, x, y + h / 2, to_x - x, to_y - (y + h / 2))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(1.5)

        callout(
            "Every resolved Discord thread ends here + in Notion",
            Inches(11.1), Inches(1.55), Inches(2.0), Inches(0.7),
            Inches(6.4), Inches(2.0),
        )
        callout(
            "24h response SLA. If it's here, it's urgent.",
            Inches(11.1), Inches(2.45), Inches(2.0), Inches(0.7),
            Inches(6.4), Inches(2.5),
        )
        callout(
            "Silent co-working. Join to signal focus mode.",
            Inches(11.1), Inches(3.35), Inches(2.0), Inches(0.7),
            Inches(6.4), Inches(3.8),
        )
    else:
        txt(s20, "[Missing Image: /mnt/user-data/uploads/1775550409325_image.png]", Inches(5.5), Inches(1.5), Inches(6), Inches(3.5), size=14, color=DANGER)

    card(s20, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), acc=DANGER)
    txt(s20, "The key discipline: when a thread resolves into a decision → #announcements + Decision Log in Notion. Discord is ephemeral. Notion is permanent.", Inches(0.8), Inches(6.6), Inches(11), Inches(0.5), size=13, color=TEXT_NAVY)
    notes(s20, "Notice the channel names include the role. Product-CTO means Arnau owns that channel. GTM-COO means Adam owns it. Finance-CFO means I own it. Nobody posts in someone else's channel without a reason.")

    # SEC 5 - 21
    s21 = add_slide()
    txt(s21, "Drive: binary blobs only", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    card(s21, Inches(1), Inches(1.5), Inches(7.5), Inches(5))
    txt(s21, "What goes in Drive", Inches(1.3), Inches(1.8), Inches(6), Inches(0.5), size=18, bold=True)
    txt(s21, "✅ PPT decks — Weekly sync presentations\n✅ Contracts & legal — NDAs, PoC agreements\n✅ Invoices & receipts — Financial paper trail\n✅ Large exports — Datasets, financial exports", Inches(1.3), Inches(2.4), Inches(7), Inches(1), size=14)
    
    rect(s21, Inches(1.3), Inches(3.5), Inches(6.9), Inches(0.02), fill=TEXT_MUTED, fl=TEXT_MUTED)
    txt(s21, "What does NOT go in Drive", Inches(1.3), Inches(3.8), Inches(6), Inches(0.5), size=18, bold=True)
    txt(s21, "❌ Meeting notes — Those go in Notion Weekly Syncs\n❌ Decisions — Those go in Notion Decision Log\n❌ Goal tracking — That's Notion Goals & OKRs", Inches(1.3), Inches(4.4), Inches(7), Inches(1), size=14)

    rect(s21, Inches(9), Inches(1.5), Inches(3.8), Inches(3), fill=RGBColor(0xF0, 0xF2, 0xF5), fl=RGBColor(0xF0, 0xF2, 0xF5))
    dtree = "Drive/\n├── Weekly Syncs/\n│   └── 2026-W15_2026-04-07/\n│       ├── oriol_w15.pptx\n│       ├── arnau_w15.pptx\n│       ├── adam_w15.pptx\n│       └── adam_notes.txt\n├── Contracts/\n├── Finance/\n└── Thesis/"
    txt(s21, dtree, Inches(9.2), Inches(1.7), Inches(3.4), Inches(2.5), size=10, name="Courier New")
    notes(s21, "If you're ever unsure whether something goes in Drive or Notion, ask: is this a file or is this knowledge? Files go in Drive. Knowledge goes in Notion.")

    # SEC 6 - 22
    add_divider(6, "Never Lose a Decision Again", "The meeting archive pipeline + AI transcription system", "This is the section I'm most excited about. Right now our meeting outputs are three PPTs in a Drive folder and whatever Adam scribbles down. We're going to fix that.")

    # 23
    s23 = add_slide()
    txt(s23, "From PPTs to permanent knowledge", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    # 5 steps
    steps = [
        ("3 PPTs + Notes", "Uploaded to Drive", RGBColor(0xDD,0xDD,0xDD), TEXT_NAVY),
        ("Claude API", "All inputs processed", TEXT_LIGHT_BLUE, TEXT_NAVY),
        ("AI Processing", "Classify, clean", ACCENT, TEXT_NAVY),
        ("Structured Output", "4 artifacts gen", SUCCESS, TEXT_NAVY),
        ("Notion", "Pasted to Sync page", DARK_BG, TEXT_WHITE)
    ]
    for i, (t1, t2, bg, fg) in enumerate(steps):
        rect(s23, Inches(0.5 + i*2.5), Inches(2), Inches(2.2), Inches(1.5), fill=bg, fl=bg)
        txt(s23, t1, Inches(0.6 + i*2.5), Inches(2.3), Inches(2), Inches(0.5), size=14, color=fg, bold=True, align=PP_ALIGN.CENTER)
        txt(s23, t2, Inches(0.6 + i*2.5), Inches(2.8), Inches(2), Inches(0.5), size=12, color=fg, align=PP_ALIGN.CENTER)
        if i < 4: txt(s23, "→", Inches(2.8 + i*2.5), Inches(2.5), Inches(0.3), Inches(0.5), size=20, bold=True)

    # 4 outputs
    for i, (t1, t2) in enumerate([("Executive summary", "5 bullets, decisions only"), ("Full writeup", "~400 words, structured"), ("Action items", "Owner + deadline + goal"), ("Decision log", "Pre-formatted for paste")]):
        rect(s23, Inches(1.5 + i*2.5), Inches(4.5), Inches(2.2), Inches(1), fill=TEXT_WHITE, fl=TEXT_NAVY)
        txt(s23, t1, Inches(1.6 + i*2.5), Inches(4.6), Inches(2), Inches(0.5), size=12, bold=True, align=PP_ALIGN.CENTER)
        txt(s23, t2, Inches(1.6 + i*2.5), Inches(5.0), Inches(2), Inches(0.5), size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Note
    txt(s23, "In 6 months: open any week, read a clean account of what was decided, who owned what, and what happened. Without wading through 3 PPTs.", Inches(0.5), Inches(6.5), Inches(12), Inches(0.5), size=13, color=TEXT_MUTED, italic=True, align=PP_ALIGN.CENTER)
    notes(s23, "Adam still takes handwritten notes in the meeting — that's important. Raw notes have context that PPTs don't. The AI takes both and produces the clean version. Adam's original notes go in Drive. The AI output goes in Notion.")

    # 24
    s24 = add_slide()
    txt(s24, "How it works — the full system", Inches(0.5), Inches(0.5), Inches(12), Inches(0.5), size=26, bold=True, align=PP_ALIGN.CENTER)
    
    # Left Flow
    y = 1.5
    for s, desc in [("1. During the meeting", "Adam takes handwritten notes on paper. No typing into Notion — stay present."), 
                    ("2. After meeting (same day)", "Adam types up raw notes as plain text. Upload to Drive with the 3 PPTs."), 
                    ("3. Oriol runs pipeline", "Feed all 4 inputs to Claude with structured prompt. ~2 minutes to run."), 
                    ("4. Paste to Notion", "Copy the 4 artifacts into the Weekly Sync Notion page. Link the Drive folder."), 
                    ("5. Future reference", "Any sync is searchable, readable, and linked to the Decision Log forever.")]:
        txt(s24, s, Inches(0.5), Inches(y), Inches(5), Inches(0.3), size=16, bold=True)
        txt(s24, desc, Inches(0.5), Inches(y + 0.3), Inches(5), Inches(0.6), size=14, color=TEXT_MUTED)
        y += 0.9

    # Right Mockup
    card(s24, Inches(6), Inches(1.5), Inches(6.8), Inches(4.5), acc=DARK_BG, bg=RGBColor(0xF0, 0xF2, 0xF5))
    mockup_text = "📅 Weekly Syncs / 2026-W15\n\nAI Summary\n  · Paver visit: qualified lead, demo scheduled Apr 14\n  · Arnau: RAG on Decfa targeted Apr 18\n  · Decision: skip Slack permanently\n\nAction Items\n  ☐ [Arnau] RAG demo on Decfa → Apr 18\n  ☐ [Adam] Paver follow-up email → Apr 8\n  ☐ [Oriol] Decision Log entries → today\n\nDecision Log links\n  → On-prem as hard constraint (reconfirmed)\n  → Slack permanently skipped\n\n📁 Drive: 2026-W15 folder"
    txt(s24, mockup_text, Inches(6.3), Inches(1.8), Inches(6), Inches(4), size=13, name="Courier New")

    card(s24, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), acc=ACCENT)
    txt(s24, "Future upgrade: Whisper local transcription (runs on our RTX 3090s) for automatic voice-to-text from The War Room. Zero cloud, fully on-prem.", Inches(0.8), Inches(6.6), Inches(11), Inches(0.5), size=13, color=TEXT_NAVY)
    notes(s24, "The future version uses Whisper running locally on our GPUs to transcribe The War Room voice calls automatically. On-prem, zero data leaves the machine — consistent with our entire product philosophy. That's a Week 9 project, not today.")

    # 25
    s25 = add_slide(DARK_BG)
    txt(s25, "Before we leave this room", Inches(0.5), Inches(0.5), Inches(12), Inches(0.8), size=36, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    def action_card(y, num, title, sub):
        card(s25, Inches(2), Inches(y), Inches(9.33), Inches(1.2))
        txt(s25, num, Inches(2.5), Inches(y + 0.3), Inches(1), Inches(0.5), size=32, color=ACCENT, bold=True)
        txt(s25, title, Inches(3.5), Inches(y + 0.2), Inches(7.5), Inches(0.5), size=20, bold=True)
        txt(s25, sub, Inches(3.5), Inches(y + 0.6), Inches(7.5), Inches(0.5), size=14, color=TEXT_MUTED)

    action_card(1.6, "01", "Fill your Q2 KRs in Notion — right now", "Goals & OKRs → Q2 2026 → your name. 3 KRs. You define them.")
    action_card(3.1, "02", "Build your W15 sprint card", "Sprint Goals → W15 → your 2–3 goals for this week with % and status.")
    action_card(4.6, "03", "Submit your first EOD log today", "Daily Logs → your name → new page. Streak starts today.")

    rect(s25, Inches(2), Inches(6.2), Inches(9.33), Inches(0.8), fill=ACCENT, fl=ACCENT)
    txt(s25, "The system is live when all three of us have a streak of 1. That happens today.", Inches(2.2), Inches(6.35), Inches(9), Inches(0.5), size=16, color=TEXT_NAVY, bold=True, align=PP_ALIGN.CENTER)
    notes(s25, "Open Notion. I'll walk you through it. We're not leaving until the streaks say 1.")

    prs.save('week15-2026.pptx')

if __name__ == '__main__':
    create_deck()
